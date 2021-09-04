from pptx import Presentation
import glob

images=glob.glob(f"./Result/{kio.PROJECT}/Fig/All/*.png")
# Presentationインスタンスの作成
ppt = Presentation("./etc/tca_theme.pptx")
# 幅
width = ppt.slide_width
# 高さ
height = ppt.slide_height
# レイアウト, 6番は白紙
blank_slide_layout = ppt.slide_layouts[6]
# ファイル毎にループ
for img in images:
    # 白紙のスライドの追加
    slide = ppt.slides.add_slide(blank_slide_layout)

    # 画像の挿入
    pic = slide.shapes.add_picture(img, 0, 0,width=width*0.8)

    # 中心に移動
    pic.left = int( ( width  - pic.width  ) / 2 )
    pic.top  = int( ( height - pic.height ) / 2 )

# 名前をつけて保存
ppt.save(f'./Result/{kio.PROJECT}/Fig/All/figure.pptx')